import os
import json
import lancedb
from sentence_transformers import SentenceTransformer
import docx
import nbformat
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

STATE_FILE_PATH = os.path.join("./lancedb_data", "index_state.json")

class RAG_System:
    def __init__(self, directory_path):
        self.directory = directory_path
        self.model = SentenceTransformer('all-mpnet-base-v2')
        self.db = lancedb.connect("./lancedb_data")
        self.table = None
        
        if "minerva_docs" in self.db.table_names():
            self.table = self.db.open_table("minerva_docs")

    def _get_folder_state(self):
        state = {}
        extensions = ['.pptx', '.docx', '.ipynb', '.txt', '.py', '.md']
        for root, _, files in os.walk(self.directory):
            for file in files:
                if file.endswith(tuple(extensions)):
                    path = os.path.join(root, file)
                    try:
                        state[path] = os.path.getmtime(path)
                    except FileNotFoundError:
                        continue # File might be deleted during scan
        return state

    def _load_index_state(self):
        if os.path.exists(STATE_FILE_PATH):
            with open(STATE_FILE_PATH, 'r') as f:
                return json.load(f)
        return None

    def _save_index_state(self, state):
        os.makedirs(os.path.dirname(STATE_FILE_PATH), exist_ok=True)
        with open(STATE_FILE_PATH, 'w') as f:
            json.dump(state, f, indent=4)

    def is_reindex_needed(self):
        if self.table is None: return True
        current_state = self._get_folder_state()
        saved_state = self._load_index_state()
        if saved_state is None or current_state != saved_state:
            return True
        return False

    def _read_file_content(self, file_path):
        try:
            if not os.path.exists(file_path): return ""
            if file_path.endswith('.pptx') and Presentation:
                prs = Presentation(file_path)
                return "\n".join(run.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame for paragraph in shape.text_frame.paragraphs for run in paragraph.runs)
            elif file_path.endswith('.docx'):
                doc = docx.Document(file_path)
                return "\n".join([para.text for para in doc.paragraphs])
            elif file_path.endswith('.ipynb'):
                with open(file_path, 'r', encoding='utf-8') as f:
                    notebook = nbformat.read(f, as_version=4)
                    return "\n\n".join(cell.source for cell in notebook.cells if cell.cell_type in ['markdown', 'code'])
            else:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
        except Exception as e:
            print(f"Error reading {file_path}: {e}")
            return ""

    def index_directory(self, status_callback):
        current_folder_state = self._get_folder_state()
        status_callback("Scanning for supported files...")
        file_paths = list(current_folder_state.keys())
        
        if not file_paths:
            status_callback("No supported files found to index.")
            return

        all_docs = []
        for i, path in enumerate(file_paths):
            status_callback(f"Reading {i+1}/{len(file_paths)}: {os.path.basename(path)}")
            content = self._read_file_content(path)
            if content:
                chunks = content.split('\n\n')
                for chunk in chunks:
                    if chunk.strip():
                        all_docs.append({"text": chunk, "source": os.path.basename(path)})
        
        if not all_docs:
            status_callback("Could not extract any text from files.")
            return
            
        status_callback(f"Creating embeddings for {len(all_docs)} text chunks...")
        embeddings = self.model.encode([doc['text'] for doc in all_docs], show_progress_bar=True)
        
        for i, doc in enumerate(all_docs):
            doc['vector'] = embeddings[i]

        status_callback("Saving embeddings to vector database...")
        if "minerva_docs" in self.db.table_names():
            self.db.drop_table("minerva_docs")
        self.table = self.db.create_table("minerva_docs", data=all_docs)
        self._save_index_state(current_folder_state)
        status_callback("Indexing complete. Ready for questions.")

    def retrieve(self, query: str, k: int = 7):
        if self.table is None: return [], []
        query_vector = self.model.encode(query)
        results = self.table.search(query_vector).limit(k).to_list()
        context_chunks = [res['text'] for res in results]
        sources = list(dict.fromkeys([res['source'] for res in results]))
        return context_chunks, sources